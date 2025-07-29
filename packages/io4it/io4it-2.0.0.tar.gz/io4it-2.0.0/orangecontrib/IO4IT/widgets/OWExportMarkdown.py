import os
import sys

from Orange.widgets import widget
from Orange.widgets.widget import Input
from AnyQt.QtWidgets import QFileDialog, QMessageBox, QApplication

from orangecontrib.text.corpus import Corpus
from docx import Document
from docx.shared import Pt as pt_docx
from pptx import Presentation
from pptx.util import Inches, Pt

if "site-packages/Orange/widgets" in os.path.dirname(os.path.abspath(__file__)).replace("\\", "/"):
    from Orange.widgets.orangecontrib.AAIT.utils.import_uic import uic
else:
    from orangecontrib.AAIT.utils.import_uic import uic

import pypandoc
import tempfile

# pip install pypandoc
# pip install pypandoc_binary

class OWExportMarkdown(widget.OWWidget):
    name = "OWExportMarkdown"
    description = "Export results from a markdown as pdf, docx, pptx"
    icon = "icons/export_md.png"
    if "site-packages/Orange/widgets" in os.path.dirname(os.path.abspath(__file__)).replace("\\", "/"):
        icon = "icons_dev/export_md.png"
    want_control_area = True
    priority = 9999

    class Inputs:
        corpus = Input("Corpus", Corpus, replaces=["Data"])

    def __init__(self):
        super().__init__()
        self.corpus = None
        self.markdown_text = ""

        # Load .ui file
        ui_path = os.path.join(os.path.dirname(__file__), "designer", "owexportmarkdown.ui")
        uic.loadUi(ui_path, baseinstance=self)

        # Connect the Save button
        self.saveButton.clicked.connect(self.save_markdown_as)

    @Inputs.corpus
    def set_corpus(self, corpus):
        if corpus is not None:
            self.corpus = corpus
            self.markdown_text = "\n".join(str(d) for d in corpus.documents)
        else:
            self.corpus = None
            self.markdown_text = ""

    def ajouter_en_tete_pied_docx(self, file_path, header_text, footer_text):
        doc = Document(file_path)
        section = doc.sections[0]

        # En-tête
        header = section.header
        if not header.paragraphs:
            p = header.add_paragraph()
        else:
            p = header.paragraphs[0]
        p.text = header_text
        p.runs[0].font.size = pt_docx(10)

        # Pied de page
        footer = section.footer
        if not footer.paragraphs:
            p = footer.add_paragraph()
        else:
            p = footer.paragraphs[0]
        p.text = footer_text
        p.runs[0].font.size = pt_docx(10)

        doc.save(file_path)

    def ajouter_entete_pied_pptx(self, file_path, entete_text, pied_text):
        prs = Presentation(file_path)
        for slide in prs.slides:
            # En-tête
            entete = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(8), Inches(0.5))
            tf_entete = entete.text_frame
            tf_entete.text = entete_text
            tf_entete.paragraphs[0].font.size = Pt(12)
            tf_entete.paragraphs[0].font.bold = True


            pied = slide.shapes.add_textbox(Inches(0.3), Inches(6.3), Inches(8), Inches(0.5))
            tf_pied = pied.text_frame
            tf_pied.text = pied_text
            tf_pied.paragraphs[0].font.size = Pt(10)

        prs.save(file_path)

    def save_markdown_as(self):
        if not self.markdown_text.strip():
            QMessageBox.warning(self, "Aucun contenu", "Aucun contenu markdown à exporter.")
            return

        filters = "PDF (*.pdf);;Word (*.docx);;PowerPoint (*.pptx)"
        file_path, selected_filter = QFileDialog.getSaveFileName(self, "Enregistrer sous", "", filters)

        if not file_path:
            return

        ext_map = {
            "PDF (*.pdf)": "pdf",
            "Word (*.docx)": "docx",
            "PowerPoint (*.pptx)": "pptx"
        }

        extension = ext_map.get(selected_filter)
        if not extension:
            QMessageBox.critical(self, "Erreur", "Format de fichier non pris en charge.")
            return

        if not file_path.lower().endswith(f".{extension}"):
            file_path += f".{extension}"

        try:
            with tempfile.NamedTemporaryFile(mode="w", suffix=".md", delete=False, encoding="utf-8") as tmp:
                tmp.write(self.markdown_text)
                tmp_path = tmp.name

            pypandoc.convert_file(tmp_path, to=extension, outputfile=file_path)

            if extension == "docx":
                self.ajouter_en_tete_pied_docx(file_path, "Rapport - Orange AI",
                                          "Page générée automatiquement - Ne pas diffuser")

            if extension == "pptx":
                self.ajouter_entete_pied_pptx(file_path, "Orange AI – Présentation", "Page générée automatiquement")

            os.remove(tmp_path)

            QMessageBox.information(self, "Succès", f"Fichier exporté :\n{file_path}")
        except Exception as e:
            QMessageBox.critical(self, "Erreur d'export", f"Erreur : {str(e)}")


if __name__ == "__main__":
    app = QApplication(sys.argv)
    my_widget = OWExportMarkdown()
    my_widget.show()
    if hasattr(app, "exec"):
        sys.exit(app.exec())
    else:
        sys.exit(app.exec_())
