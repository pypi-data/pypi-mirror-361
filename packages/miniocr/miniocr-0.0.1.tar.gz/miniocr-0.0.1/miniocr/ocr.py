import asyncio
import base64
import os
from typing import List, Optional
from openai import AsyncOpenAI
import aiofiles
import aiohttp
from pdf2image import convert_from_path
import tempfile
from pptx import Presentation

class MiniOCR:
    def __init__(self, api_key: str = None):
        self.client = AsyncOpenAI(api_key=api_key or os.getenv("OPENAI_API_KEY"))
        self.system_prompt = """
Convert the following document to markdown.
Return only the markdown with no explanation text. Do not include delimiters like ```markdown or ```html.

RULES:
  - You must include all information on the page. Do not exclude headers, footers, or subtext.
  - Return tables in an HTML format.
  - Charts & infographics must be interpreted to a markdown format. Prefer table format when applicable.
  - Logos should be wrapped in brackets. Ex: <logo>Coca-Cola<logo>
  - Watermarks should be wrapped in brackets. Ex: <watermark>OFFICIAL COPY<watermark>
  - Page numbers should be wrapped in brackets. Ex: <page_number>14<page_number> or <page_number>9/22<page_number>
  - Prefer using ☐ and ☑ for check boxes.
"""

    async def encode_image_to_base64(self, image_path: str) -> str:
        """Encode an image to base64 asynchronously."""
        async with aiofiles.open(image_path, "rb") as image_file:
            image_data = await image_file.read()
        return base64.b64encode(image_data).decode("utf-8")

    def is_image_file(self, file_path: str) -> bool:
        """Check if file is an image."""
        image_extensions = ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp')
        return file_path.lower().endswith(image_extensions)

    def is_pdf_file(self, file_path: str) -> bool:
        """Check if file is a PDF."""
        return file_path.lower().endswith('.pdf')

    def is_pptx_file(self, file_path: str) -> bool:
        """Check if file is a PPTX."""
        return file_path.lower().endswith('.pptx')

    async def download_file(self, url: str, temp_dir: str) -> str:
        """Download file from URL if needed."""
        if url.startswith(('http://', 'https://')):
            async with aiohttp.ClientSession() as session:
                async with session.get(url) as response:
                    filename = os.path.basename(url.split('?')[0]) or 'document'
                    filepath = os.path.join(temp_dir, filename)
                    async with aiofiles.open(filepath, 'wb') as f:
                        async for chunk in response.content.iter_chunked(8192):
                            await f.write(chunk)
                    return filepath
        return url

    async def pdf_to_images(self, pdf_path: str, temp_dir: str, dpi: int = 200) -> List[str]:
        """Convert PDF to images."""
        images = convert_from_path(pdf_path, dpi=dpi)
        image_paths = []
        
        for i, image in enumerate(images):
            image_path = os.path.join(temp_dir, f"page_{i+1}.png")
            image.save(image_path, 'PNG')
            image_paths.append(image_path)
        
        return image_paths

    async def pptx_to_images_and_text(self, pptx_path: str, temp_dir: str) -> (List[str], str):
        """Extract images and text from a PPTX file."""
        prs = Presentation(pptx_path)
        image_paths = []
        text_content = []
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image = shape.image
                    image_bytes = image.blob
                    image_filename = f"slide_{i+1}_{shape.name}.{image.ext}"
                    image_path = os.path.join(temp_dir, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    image_paths.append(image_path)
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return image_paths, "\n".join(text_content)

    async def process_image(self, image_path: str, model: str = "gpt-4o-mini") -> str:
        """Process a single image with OpenAI Vision API."""
        base64_image = await self.encode_image_to_base64(image_path)
        
        response = await self.client.chat.completions.create(
            model=model,
            messages=[
                {
                    "role": "system",
                    "content": self.system_prompt
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=4000
        )
        
        return response.choices[0].message.content

    async def ocr(
        self,
        file_path: str,
        model: str = "gpt-4o-mini",
        concurrency: int = 5,
        output_dir: Optional[str] = None,
        cleanup: bool = True
    ) -> dict:
        """
        Main function to convert PDF/image/pptx to markdown using OpenAI Vision API.
        
        Args:
            file_path: Path or URL to PDF or image file
            model: OpenAI model to use (default: gpt-4o-mini)
            concurrency: Number of concurrent requests
            output_dir: Directory to save markdown output
            cleanup: Whether to cleanup temporary files
            
        Returns:
            Dictionary with markdown content and metadata
        """
        with tempfile.TemporaryDirectory() as temp_dir:
            text_content = ""
            # Download file if URL
            local_path = await self.download_file(file_path, temp_dir)
            
            # Determine file type and get image paths
            if self.is_image_file(local_path):
                image_paths = [local_path]
            elif self.is_pdf_file(local_path):
                image_paths = await self.pdf_to_images(local_path, temp_dir)
            elif self.is_pptx_file(local_path):
                image_paths, text_content = await self.pptx_to_images_and_text(local_path, temp_dir)
            else:
                raise ValueError(f"Unsupported file type: {local_path}")
            
            # Process images concurrently
            semaphore = asyncio.Semaphore(concurrency)
            
            async def process_with_semaphore(image_path):
                async with semaphore:
                    return await self.process_image(image_path, model)
            
            tasks = [process_with_semaphore(img_path) for img_path in image_paths]
            results = await asyncio.gather(*tasks)
            
            # Combine results
            if not results:
                markdown_content = text_content
            else:
                markdown_content = "\n\n".join(results) if len(results) > 1 else results[0]
                if text_content:
                    markdown_content = f"{text_content}\n\n{markdown_content}"
            
            # Save to file if output_dir specified
            if output_dir:
                os.makedirs(output_dir, exist_ok=True)
                filename = os.path.splitext(os.path.basename(local_path))[0]
                output_path = os.path.join(output_dir, f"{filename}.md")
                async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
                    await f.write(markdown_content)
            
            return {
                "content": markdown_content,
                "pages": len(results),
                "file_name": os.path.basename(local_path)
            }