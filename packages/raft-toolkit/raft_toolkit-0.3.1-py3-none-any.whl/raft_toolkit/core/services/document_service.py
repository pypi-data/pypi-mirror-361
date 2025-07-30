"""
Document processing service for chunking and text extraction.
"""

import json
import logging
import re
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from math import ceil
from pathlib import Path
from typing import Any, List

from ..config import RaftConfig
from ..models import DocumentChunk
from .embedding_service import create_embedding_service
from .llm_service import LLMService

try:
    import pypdf
except ImportError:
    pypdf = None  # type: ignore

try:
    from pptx import Presentation

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

    class Presentation:  # type: ignore[no-redef]
        def __init__(self, *args: Any, **kwargs: Any) -> None:
            pass


try:
    from langchain_openai.embeddings import AzureOpenAIEmbeddings, OpenAIEmbeddings

    HAS_LANGCHAIN_OPENAI = True
except ImportError:
    HAS_LANGCHAIN_OPENAI = False

    class OpenAIEmbeddings:  # type: ignore[no-redef]
        def __init__(self, *args: Any, **kwargs: Any) -> None:
            pass

    class AzureOpenAIEmbeddings:  # type: ignore[no-redef]
        def __init__(self, *args: Any, **kwargs: Any) -> None:
            pass


try:
    from langchain_community.embeddings import NomicEmbeddings

    HAS_NOMIC_EMBEDDINGS = True
except ImportError:
    HAS_NOMIC_EMBEDDINGS = False

    class NomicEmbeddings:  # type: ignore[no-redef]
        def __init__(self, *args: Any, **kwargs: Any) -> None:
            pass


try:
    from langchain_experimental.text_splitter import SemanticChunker

    HAS_SEMANTIC_CHUNKER = True
except ImportError:
    HAS_SEMANTIC_CHUNKER = False

    class SemanticChunker:  # type: ignore[no-redef]
        def __init__(self, *args: Any, **kwargs: Any) -> None:
            pass


# Overall LangChain availability
HAS_LANGCHAIN = HAS_LANGCHAIN_OPENAI or HAS_NOMIC_EMBEDDINGS


try:
    from tqdm import tqdm

    HAS_TQDM = True
except ImportError:
    HAS_TQDM = False

    # Create a fallback tqdm function if not available
    def tqdm(iterable: Any, *args: Any, **kwargs: Any) -> Any:  # type: ignore[no-redef]
        return iterable


logger = logging.getLogger(__name__)


class DocumentService:
    """Service for document processing and chunking."""

    def __init__(self, config: RaftConfig, llm_service: LLMService):
        self.config = config
        self.llm_service = llm_service
        self.embedding_service = create_embedding_service(config)

    def process_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process documents and return chunks."""
        logger.info(f"Processing documents from {data_path}")

        if self.config.doctype == "api":
            return self._process_api_documents(data_path)
        else:
            return self._process_regular_documents(data_path)

    def _process_api_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process API documentation from JSON file."""
        with open(data_path) as f:
            api_docs_json = json.load(f)

        required_fields = ["user_name", "api_name", "api_call", "api_version", "api_arguments", "functionality"]
        if api_docs_json and isinstance(api_docs_json[0], dict):
            for field in required_fields:
                if field not in api_docs_json[0]:
                    raise ValueError(f"API documentation missing required field: {field}")

        chunks = []
        for i, api_doc in enumerate(api_docs_json):
            chunk = DocumentChunk.create(
                content=str(api_doc), source=str(data_path), metadata={"type": "api", "index": i}
            )
            chunks.append(chunk)

        return chunks

    def _process_regular_documents(self, data_path: Path) -> List[DocumentChunk]:
        """Process regular documents (PDF, TXT, JSON, PPTX)."""
        embeddings = self._build_embeddings()

        # Get list of files to process
        file_paths = []
        if data_path.is_dir():
            file_paths = list(data_path.rglob(f"**/*.{self.config.doctype}"))
        else:
            file_paths = [data_path]

        all_chunks = []
        futures = []

        with tqdm(total=len(file_paths), desc="Processing files", unit="file") as pbar:
            with ThreadPoolExecutor(max_workers=self.config.embed_workers) as executor:
                for file_path in file_paths:
                    future = executor.submit(self._process_single_file, embeddings, file_path)
                    futures.append(future)

                    if self.config.pace:
                        time.sleep(15)

                for future in as_completed(futures):
                    try:
                        chunks = future.result()
                        all_chunks.extend(chunks)
                        pbar.set_postfix({"total_chunks": len(all_chunks)})
                        pbar.update(1)
                    except Exception as e:
                        logger.error(f"Error processing file: {e}")
                        pbar.update(1)

        # Create embeddings using the embedding service with template
        if self.config.embedding_prompt_template or True:  # Always use the embedding service for consistency
            all_chunks = self.embedding_service.create_embeddings_with_template(all_chunks)

        return all_chunks

    def _process_single_file(self, embeddings: Any, file_path: Path) -> List[DocumentChunk]:
        """Process a single file and return its chunks."""
        logger.debug(f"Processing file: {file_path}")

        # Extract text based on document type
        text = self._extract_text(file_path)

        # Split into chunks
        chunk_contents = self._split_text(embeddings, text)

        # Create DocumentChunk objects
        chunks = []
        for i, content in enumerate(chunk_contents):
            chunk = DocumentChunk.create(
                content=content,
                source=str(file_path),
                metadata={
                    "type": self.config.doctype,
                    "chunk_index": i,
                    "chunking_strategy": self.config.chunking_strategy,
                },
            )
            chunks.append(chunk)

        return chunks

    def _extract_text(self, file_path: Path) -> str:
        """Extract text from a file based on its type."""
        if self.config.doctype == "json":
            with open(file_path, "r") as f:
                data = json.load(f)
            text_value = data.get("text", str(data))
            return str(text_value)  # Ensure we return a string

        elif self.config.doctype == "pdf":
            text = ""
            with open(file_path, "rb") as file:
                reader = pypdf.PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text()
            return text

        elif self.config.doctype == "txt":
            with open(file_path, "r", encoding="utf-8") as file:
                return file.read()

        elif self.config.doctype == "pptx":
            return self._extract_text_from_pptx(file_path)

        else:
            raise ValueError(f"Unsupported document type: {self.config.doctype}")

    def _extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint file."""
        # Convert Path to string for compatibility with Presentation
        file_path_str = str(file_path)
        prs = Presentation(file_path_str)
        text_parts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
                elif hasattr(shape, "table"):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text_parts.append(cell.text)

        return "\n".join(text_parts)

    def _split_text(self, embeddings: Any, text: str) -> List[str]:
        """Split text into chunks based on the configured strategy."""
        if self.config.chunking_strategy == "semantic":
            return self._semantic_chunking(embeddings, text)
        elif self.config.chunking_strategy == "fixed":
            return self._fixed_chunking(text)
        elif self.config.chunking_strategy == "sentence":
            return self._sentence_chunking(text)
        else:
            raise ValueError(f"Unknown chunking strategy: {self.config.chunking_strategy}")

    def _semantic_chunking(self, embeddings: Any, text: str) -> List[str]:
        """Perform semantic chunking using embeddings."""
        if not HAS_SEMANTIC_CHUNKER:
            logger.warning("SemanticChunker not available, falling back to fixed chunking")
            return self._fixed_chunking(text)

        try:
            params = self.config.chunking_params
            num_chunks = params.get("number_of_chunks") or ceil(len(text) / self.config.chunk_size)
            min_chunk_size = params.get("min_chunk_size", 0)

            # Ensure we have a reasonable number of chunks
            if num_chunks <= 0:
                num_chunks = 1

            # Create semantic chunker with correct parameters
            text_splitter = SemanticChunker(
                embeddings=embeddings, number_of_chunks=num_chunks, min_chunk_size=min_chunk_size
            )

            chunks = text_splitter.create_documents([text])

            # Validate chunks
            if not chunks:
                logger.warning("Semantic chunking produced no chunks, falling back to fixed chunking")
                return self._fixed_chunking(text)

            return [chunk.page_content for chunk in chunks]
        except Exception as e:
            logger.error(f"Error during semantic chunking: {e}, falling back to fixed chunking")
            return self._fixed_chunking(text)

    def _fixed_chunking(self, text: str) -> List[str]:
        """Split text into fixed-size chunks."""
        chunk_size = self.config.chunk_size
        return [text[i : i + chunk_size] for i in range(0, len(text), chunk_size)]

    def _sentence_chunking(self, text: str) -> List[str]:
        """Split text by sentences, respecting chunk size limits."""
        sentences = re.split(r"(?<=[.!?]) +", text)
        chunks = []
        current_chunk = ""

        for sentence in sentences:
            if len(current_chunk) + len(sentence) > self.config.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = sentence
            else:
                current_chunk += (" " if current_chunk else "") + sentence

        if current_chunk:
            chunks.append(current_chunk)

        return chunks

    def _build_embeddings(self) -> Any:
        """Build embeddings model for semantic chunking."""
        try:
            # First try to use the client builder function
            try:
                from ..clients.openai_client import build_langchain_embeddings

                if self.config.use_azure_identity:
                    from ..utils import get_azure_openai_token

                    api_key = get_azure_openai_token()
                else:
                    api_key = self.config.openai_key

                return build_langchain_embeddings(api_key=api_key, model=self.config.embedding_model)
            except ImportError:
                # If client builder is not available, try direct initialization

                # Check if we should use Nomic embeddings
                if self.config.embedding_model.startswith("nomic-"):
                    if HAS_NOMIC_EMBEDDINGS:
                        logger.info(f"Using Nomic embeddings model: {self.config.embedding_model}")
                        return NomicEmbeddings(model=self.config.embedding_model)
                    else:
                        logger.warning("Nomic embeddings requested but not available")

                # Fall back to OpenAI embeddings
                if self.config.azure_openai_enabled:
                    if not HAS_LANGCHAIN_OPENAI:
                        return self._create_mock_embeddings()

                    # Parameters for Azure OpenAI
                    kwargs = {
                        "deployment": self.config.embedding_model,
                        "api_version": "2023-05-15",
                    }

                    if hasattr(self.config, "azure_endpoint"):
                        kwargs["endpoint"] = self.config.azure_endpoint

                    if self.config.use_azure_identity:
                        from ..utils import get_azure_openai_token

                        api_key = get_azure_openai_token()
                    else:
                        api_key = self.config.openai_key

                    if api_key:
                        kwargs["api_key"] = api_key

                    return AzureOpenAIEmbeddings(**kwargs)  # type: ignore
                else:
                    if not HAS_LANGCHAIN_OPENAI:
                        return self._create_mock_embeddings()

                    # Parameters for OpenAI
                    kwargs = {}

                    if self.config.embedding_model:
                        kwargs["model"] = self.config.embedding_model

                    if self.config.openai_key:
                        kwargs["api_key"] = self.config.openai_key

                    return OpenAIEmbeddings(**kwargs)  # type: ignore
        except Exception as e:
            logger.error(f"Error building embeddings model: {e}")
            return self._create_mock_embeddings()

    def _create_mock_embeddings(self) -> Any:
        """Create a mock embeddings model for testing or when real implementation is unavailable."""

        class MockEmbeddings:
            def embed_documents(self, texts: List[str]) -> List[List[float]]:
                return [[0.1, 0.2, 0.3] for _ in texts]

            def embed_query(self, text: str) -> List[float]:
                return [0.1, 0.2, 0.3]

        logger.warning("Using mock embeddings implementation")
        return MockEmbeddings()
