import json
import os
from colorama import Fore, Back, Style
from datetime import datetime
import re
import traceback
import base64
import PyPDF2
import docx
from pptx import Presentation
import pandas as pd
from PIL import Image
import requests
import mimetypes
from urllib.parse import urlparse, unquote
from hero.util import log
from typing import List, Dict
import mimetypes
import chardet
from ulid import ULID
from bs4 import BeautifulSoup, Comment
from playwright.sync_api import Page


# TODO: 后期优化，在 Hero 中可以自定义这些参数
SPLIT_FILE_LIMIT = 50000
LINE_LIMIT = 100000
MAX_HISTORY_COUNT = 20
MAX_HISTORY_LIMIT = 40000

def read_default_tool_prompt(path: str) -> str:
    """
    读取默认提示词
    """
    # 读取相对路径 ../prompt/xxx.md
    current_dir = os.path.dirname(__file__)
    with open(os.path.join(current_dir, "..", "prompt", path), "r", encoding="utf-8") as f:
        return f.read()

def extract_tool_response(content):
    """从内容中提取工具名称和参数，根据简单的JSON格式"""
    try:
        pattern = r"```json[\s\S]*```"
        match = re.search(pattern, content)
        if match:
            content = match.group().replace("\xa0", " ")
            content = content.replace("```json", "").replace("```", "")
            content = content.strip()
            return json.loads(content)
        else:
            return None

    except Exception as e:
        log.error(f"extract_tool_response error: {e}")
        log.error(traceback.format_exc())
        return None


def write_task_execute_history(tool, params, status, message, caller):
    """
    写入任务执行历史
    """
    with open(
        os.path.join(caller.get("log_dir"), f"__{caller.get('name')}_tasks.md"),
        "a",
        encoding="utf-8",
    ) as f:
        content = ""
        content += f"<task tool={tool} index={caller.get('index')} finished_at={timestamp_to_str(timestamp())}>\n"
        content += f"<params>{json.dumps(params, ensure_ascii=False)}</params>\n"
        content += f"<result status={status}>\n"
        content += message
        content += f"\n</result>\n"
        content += f"</task>\n\n"
        f.write(content)


def read_task_execute_history(caller):
    """
    读取任务执行历史
    """
    # 如果文件不存在，则创建
    file_path = os.path.join(caller.get("log_dir"), f"__{caller.get('name')}_tasks.md")
    if not os.path.exists(file_path):
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("")

    start_index = 1
    total_count = 0
    history = ""
    with open(file_path, "r", encoding="utf-8") as f:
        history = f.read()

    history_list = history.split("</task>")
    total_count = len(history_list)
    log.debug(f"len(history_list): {len(history_list)}")
    log.debug(f"len(history): {len(history)}")
    if len(history_list) > MAX_HISTORY_COUNT:
        # 取最后 MAX_HISTORY_COUNT 个
        # 计算 start_index，是整体任务列表的索引
        start_index = len(history_list) - MAX_HISTORY_COUNT + 1
        history_list = history_list[-MAX_HISTORY_COUNT:]

    # 把 history_list 拼接成字符串，如果 history_list 的字符串长度大于 MAX_HISTORY_LIMIT，则从后往前截取，直到小于 MAX_HISTORY_LIMIT
    history = "</task>".join(history_list)
    while len(history) > MAX_HISTORY_LIMIT:
        history_list.pop(0)
        start_index += 1
        history = "</task>".join(history_list)

    history = "</task>".join(history_list)
    log.debug(f"After len(history_list): {len(history_list)}")
    log.debug(f"After len(history): {len(history)}")
    log.debug(f"start_index: {start_index}")

    return history, start_index, total_count


def read_user_message(caller):
    """
    读取用户消息
    """
    file_path = os.path.join(caller.get("log_dir"), "__meta.json")
    if not os.path.exists(file_path):
        return ""

    with open(file_path, "r", encoding="utf-8") as f:
        meta = json.loads(f.read())
        return meta.get("user_message")


def write_user_message(caller, message):
    """
    写入用户消息
    """
    file_path = os.path.join(caller.get("log_dir"), "__meta.json")

    meta = json.loads(read_file(caller.get("log_dir"), "__meta.json") or "{}")

    if not meta.get("user_message"):
        meta["user_message"] = [message]
    else:
        meta["user_message"].append(message)

    with open(file_path, "w", encoding="utf-8") as f:
        f.write(json.dumps(meta, ensure_ascii=False, indent=4))

    with open(
        os.path.join(caller.get("log_dir"), f"__{caller.get('name')}_tasks.md"),
        "a",
        encoding="utf-8",
    ) as f:
        content = ""
        content += f"<user_message>\n"
        content += message
        content += f"\n</user_message>\n\n"
        f.write(content)

def backtrack_task_execute_history(caller, index):
    file_path = os.path.join(caller.get("log_dir"), f"__{caller.get('name')}_tasks.md")
    with open(file_path, "r", encoding="utf-8") as f:
        history = f.read()
    task_pattern = re.compile(r"(<task[^>]*index=(\d+)[^>]*>.*?</task>)", re.DOTALL)
    pos = 0
    new_history = ""
    for m in task_pattern.finditer(history):
        task_index = int(m.group(2))
        if task_index == index:
            new_history += history[pos:m.start()]
            break
        new_history += history[pos:m.end()]
        pos = m.end()
    else:
        # 可能有task后还有内容（如空行），保留到当前pos
        new_history += history[pos:pos]
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(new_history)

def write_ans(caller, data):
    """
    写入答案
    """
    meta = json.loads(read_file(caller.get("log_dir"), "__meta.json") or "{}")
    meta["answer"] = data
    write_file(caller.get("log_dir"), "__meta.json", json.dumps(meta, ensure_ascii=False, indent=4))

def write_file(dir, file_name, content):
    """
    写入文件
    """
    file_path = os.path.join(dir, file_name)
    if not os.path.exists(dir):
        os.makedirs(dir)
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(content)

def read_file(dir, file_name):
    """
    读取文件
    """
    file_path = os.path.join(dir, file_name)
    if not os.path.exists(file_path):
        return ""
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def append_file(dir, file_name, content):
    file_path = os.path.join(dir, file_name)
    with open(file_path, "a", encoding="utf-8") as f:
        f.write(content)

def ulid():
    """
    生成ULID
    """
    return str(ULID())

def now():
    """
    获取当前时间
    """
    return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")


def timestamp():
    """
    获取当前时间戳
    """
    return str(int(datetime.now().timestamp()))


def timestamp_to_str(timestamp: str):
    """
    将时间戳转换为字符串
    """
    return datetime.fromtimestamp(int(timestamp)).strftime("%H:%M:%S")


def print_progress(content, timestamp: str):
    """
    打印进度
    """
    print(
        f"\n{Fore.CYAN}## {content.upper()} {Style.RESET_ALL} - {timestamp_to_str(timestamp)}"
    )


def print_block(content, timestamp: str):
    """
    打印块
    """
    print(
        f"\n{Back.BLUE}## {content.upper()} {Style.RESET_ALL} - {timestamp_to_str(timestamp)}"
    )


def print_kv(key, value):
    """
    打印键值对
    """
    print(
        f"{Fore.LIGHTBLUE_EX}[{Style.RESET_ALL}{key}{Fore.LIGHTBLUE_EX}]{Style.RESET_ALL}: {Fore.LIGHTBLUE_EX}{value}{Style.RESET_ALL}"
    )

def image_to_base64_url(image_path):
    """
    将图片转换为base64编码的data URL
    """
    try:
        # 打开图片并获取格式
        with Image.open(image_path) as img:
            img_format = img.format  # 如 'JPEG', 'PNG', 'GIF' 等

            if img_format is None:
                img_format = "JPEG"

            # 将图片转为 base64 编码
            with open(image_path, "rb") as image_file:
                img_base64 = base64.b64encode(image_file.read()).decode("utf-8")

            # 构建 data URL
            mime_type = f"image/{img_format.lower()}"
            if img_format.upper() == "JPG":  # 有些系统用 'JPG'，需要转换为 'jpeg'
                mime_type = "image/jpeg"

            log.debug(f"img_base64 length: {len(img_base64)}")
            log.debug(f"mime_type: {mime_type}")

            base64_url = f"data:{mime_type};base64,{img_base64}"

            return {
                "status": "success",
                "base64_url": base64_url,
            }

    except Exception as e:
        log.error(f"将图片转换为base64编码的data URL时出错: {e}")
        log.error(traceback.format_exc())
        return {
            "status": "error",
            "message": str(e),
        }


def transfer_pdf_to_markdown(file_path: str) -> str:
    """
    将PDF文件转换为Markdown格式
    """
    try:
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            markdown_content = []

            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text.strip():
                    markdown_content.append(f"## 第 {page_num + 1} 页\n\n{text}\n")

            return "\n".join(markdown_content)
    except Exception as e:
        log.error(f"转换PDF到Markdown时出错: {e}")
        return f"转换PDF到Markdown时出错: {str(e)}"


def transfer_docx_to_markdown(file_path: str) -> str:
    """
    将DOCX文件转换为Markdown格式
    """
    try:
        doc = docx.Document(file_path)
        markdown_content = []

        for para in doc.paragraphs:
            if para.text.strip():
                # 处理标题
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    level = int(para.style.name[-1])
                    markdown_content.append(f"{'#' * level} {para.text}\n")
                else:
                    markdown_content.append(f"{para.text}\n")

        # 处理表格
        for table in doc.tables:
            markdown_table = []
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                markdown_table.append("| " + " | ".join(cells) + " |")
            if len(markdown_table) > 0:
                markdown_table.insert(
                    1, "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                )
                markdown_content.append("\n".join(markdown_table) + "\n")

        return "\n".join(markdown_content)
    except Exception as e:
        log.error(f"转换DOCX到Markdown时出错: {e}")
        return f"转换DOCX到Markdown时出错: {str(e)}"


def transfer_pptx_to_markdown(file_path: str) -> str:
    """
    将PPTX文件转换为Markdown格式
    """
    try:
        prs = Presentation(file_path)
        markdown_content = []

        for i, slide in enumerate(prs.slides):
            markdown_content.append(f"## 幻灯片 {i + 1}\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip(): # type: ignore
                    markdown_content.append(f"{shape.text}\n") # type: ignore
            markdown_content.append("\n---\n")

        return "\n".join(markdown_content)
    except Exception as e:
        log.error(f"转换PPTX到Markdown时出错: {e}")
        return f"转换PPTX到Markdown时出错: {str(e)}"


def transfer_xlsx_to_markdown(file_path: str) -> str:
    """
    将XLSX文件转换为Markdown格式
    """
    try:
        excel_file = pd.ExcelFile(file_path)
        markdown_content = []

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            markdown_content.append(f"## {sheet_name}\n")

            # 将DataFrame转换为Markdown表格
            markdown_table = df.to_markdown(index=False)
            markdown_content.append(markdown_table)
            markdown_content.append("\n")

        return "\n".join(markdown_content)
    except Exception as e:
        log.error(f"转换XLSX到Markdown时出错: {e}")
        return f"转换XLSX到Markdown时出错: {str(e)}"


def list_files_recursive(dir: str) -> str:
    """递归展示目录下的所有文件，以 Markdown 格式返回

    Args:
        workspace: 工作空间名称

    Returns:
        str: Markdown 格式的文件列表
    """
    if not os.path.exists(dir):
        return "目录不存在"

    def _list_files_recursive(path: str, level: int = 0) -> str:
        content = ""
        try:
            entries = os.listdir(path)
            for entry in sorted(entries):
                if entry.startswith("__"):
                    continue

                full_path = os.path.join(path, entry)
                relative_path = os.path.relpath(full_path, dir)

                if os.path.isdir(full_path):
                    content += f"{'  ' * level}- {entry}/\n"
                    content += _list_files_recursive(full_path, level + 1)
                else:
                    content += f"{'  ' * level}- {entry}\n"
        except Exception as e:
            log.error(f"Error listing files in {path}: {e}")
            content += f"{'  ' * level}- Error accessing directory\n"

        return content
        
    return _list_files_recursive(dir)

def is_downloadable(url: str) -> bool:
    """
    判断URL是否可下载
    """
    if url.lower().endswith(
        (
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".webp",
            ".pdf",
            ".doc",
            ".docx",
            ".xls",
            ".xlsx",
            ".ppt",
            ".pptx",
        )
    ):
        return True
    
    try:
        response = requests.head(url, timeout=10, allow_redirects=True)
        content_type = response.headers.get("content-type", "").lower()
        log.debug(f"content_type: {content_type}")
        if "text/html" in content_type:
            return False
        else:
            return True

    except Exception as e:
        return False


def get_filename_from_cd(content_disposition):
    """
    从 Content-Disposition 头部提取文件名。
    """
    if not content_disposition:
        return None
    # 尝试匹配 filename*=
    fname = re.findall("filename\\*=([^;]+)", content_disposition)
    if fname:
        encoding_filename = fname[0].split("''")
        if len(encoding_filename) == 2:
            filename = unquote(encoding_filename[1])
            return filename
    # 尝试匹配 filename=
    fname = re.findall('filename="?([^";]+)"?', content_disposition)
    if fname:
        return fname[0]
    return None


def download_file(url, save_dir=".", timeout=30):
    """
    下载文件
    """
    try:
        # 发送 GET 请求
        response = requests.get(url, stream=True, timeout=timeout)
        response.raise_for_status()

        content_type = response.headers.get("Content-Type", "")
        log.debug(f"content_type: {content_type}")

        if "text/html" in content_type:
            return False

        content_type = response.headers.get('Content-Type')
        log.debug(f"content_type: {content_type}")

        if content_type and "text/html" in content_type:
            return False

        # 从 Content-Disposition 中获取文件名
        content_disposition = response.headers.get("Content-Disposition")
        filename = get_filename_from_cd(content_disposition)

        # 如果无法从 Content-Disposition 获取文件名，则从 URL 中提取
        if not filename:
            parsed_url = urlparse(url)
            filename = os.path.basename(parsed_url.path)
            if not filename:
                # 如果 URL 中也无法获取文件名，则使用默认名称
                filename = "downloaded_file"
                # 尝试根据 Content-Type 添加扩展名
                content_type = response.headers.get("Content-Type")
                if content_type:
                    ext = mimetypes.guess_extension(content_type.split(";")[0].strip())
                    if ext:
                        filename += ext

        # 创建保存路径
        save_path = os.path.join(save_dir, filename)

        # 保存文件
        with open(save_path, "wb") as f:
            for chunk in response.iter_content(chunk_size=8192):
                if chunk:
                    f.write(chunk)

        print(f"文件已保存至: {save_path}")
        return filename

    except requests.RequestException as e:
        print(f"下载失败: {e}")
        return None


def is_binary_file(file_path, blocksize=1024):
    """
    判断一个文件是否为二进制文件。
    优先根据文件内容判断，必要时结合扩展名和 MIME 类型作为辅助。

    参数:
        file_path: 文件路径
        blocksize: 读取文件的字节数（默认读取前1024字节）

    返回:
        True 表示可能是二进制文件，False 表示是文本文件
    """
    # 先通过后缀判断一些常见的纯文本
    plain_text_extensions = {
        '.txt', '.csv', '.log', '.md', '.xml', '.json', 
        '.ini', '.cfg', '.properties', '.yaml', '.yml', 
        '.tsv', '.html', '.sql', '.bat', '.sh', 
        '.py', '.pl', '.rb', '.go', '.r', 
        '.tex', '.diff', '.patch', '.svg', '.markdown',
        '.cpp', '.c', '.h', '.java', '.js', 
        '.ts', '.swift', '.php', '.css', '.scss', 
        '.vb', '.asm', '.sql', '.dart', '.kotlin',
        '.rs', '.clj', '.cljs', '.lua', '.pl', 
        '.erl', '.ex', '.exs', '.f90', '.f95'
    }
    
    # 获取文件的后缀
    ext = os.path.splitext(file_path)[1].lower()
    if ext in plain_text_extensions:
        return False

    # 方法一：读取部分内容检测不可打印字符
    try:
        
        with open(file_path, "rb", encoding="utf-8") as file:
            all_file_content = file.read()
            chunk = all_file_content[:blocksize]
            if not chunk:
                print("空文件")
                return False  # 空文件当作文本处理
            if b"\0" in chunk:
                print("包含 NULL 字节")
                return True  # 包含 NULL 字节，大概率是二进制
            
            # 使用 chardet 检测编码，目的是兼容中文
            result = chardet.detect(chunk)
            encoding = result['encoding']
            confidence = result['confidence']
            # 尝试解码
            if confidence > 0.5:  # 只在置信度较高时解码
                try:
                    all_file_content.decode(encoding)
                    return False  # 如果成功解码，认为是文本文件
                except (UnicodeDecodeError, TypeError):
                    pass  # 解码失败，继续检查其他条件

            # 定义文本字符范围
            text_chars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x7F)))
            nontext = chunk.translate(None, text_chars)

            if float(len(nontext)) / len(chunk) > 0.8:
                return True
    except Exception as e:
        print(f"读取文件时出错: {e}")
        return True  # 保守起见，无法读取时视为二进制

    # 方法二：根据 MIME 类型判断（辅助）
    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type is not None and mime_type.startswith('text'):
        return False
    elif mime_type is not None:
        return True

    # 方法三：根据文件扩展名判断（辅助）
    binary_exts = {
        '.exe', '.bin', '.dat', '.jpg', '.jpeg', '.png', '.gif', '.pdf',
        '.zip', '.tar', '.gz', '.7z', '.mp3', '.mp4', '.avi', '.mov',
    }
    ext = os.path.splitext(file_path)[1].lower()
    if ext in binary_exts:
        return True

    # 默认认为是文本
    return False

def split_markdown(content: str) -> List[str]:
    """
    将 Markdown 内容按照标题分割，并确保每部分不超过指定长度
    """
    blocks = content.split("## ")
    temp_str = ""
    text_array = []

    for i in range(len(blocks)):
        if len(temp_str + blocks[i]) > SPLIT_FILE_LIMIT:
            split_array = split_line_by_length(temp_str, SPLIT_FILE_LIMIT)
            text_array.extend(split_array)
            temp_str = ""
        temp_str += "## " + blocks[i]

    split_array = split_line_by_length(temp_str, SPLIT_FILE_LIMIT)
    text_array.extend(split_array)
    return text_array


def split_line_by_length(text: str, length: int) -> List[str]:
    """
    将文本按行分割，确保每一段不超过指定长度
    """
    lines = text.split("\n")
    result = []
    temp_str = ""

    for i in range(len(lines)):
        if len(temp_str) > length:
            split_array = split_string_by_length(temp_str, length)
            result.extend(split_array)
            temp_str = ""
        if len(temp_str + lines[i]) > length:
            result.append(temp_str)
            temp_str = ""
        temp_str += lines[i]

    if len(temp_str) > SPLIT_FILE_LIMIT:
        split_array = split_string_by_length(temp_str, length)
        for part in split_array:
            print(len(part))
        result.extend(split_array)
    elif len(temp_str) > 0:
        result.append(temp_str)

    return result


def split_string_by_length(text: str, length: int) -> List[str]:
    """
    将字符串按指定长度分割
    """
    return [text[i : i + length] for i in range(0, len(text), length)]


def read_file_n_lines(dir: str, file_name: str, n: int) -> str:
    """
    读取文件的前n行
    """
    file_path = os.path.join(dir, file_name)
    if not os.path.exists(file_path):
        return f"File {file_name} not found"

    try:
        if file_name.lower().endswith(".xlsx") or file_name.lower().endswith(".xls"):
            # Excel文件
            df = pd.read_excel(file_path)
            content = (
                f"<!-- Preview {n} rows of {file_name}, total {len(df)} rows -->\n\n"
            )
            # 转成字符串
            content += df.head(n).to_csv(index=False)
        else:
            # 普通文本文件
            with open(file_path, "r", encoding="utf-8") as f:
                lines = f.readlines()
            content = (
                f"<!-- Preview {n} lines of {file_name}, total {len(lines)} lines -->\n\n"
            )
            for line in lines[:n]:
                content += line
        return content
    except Exception as e:
        log.error(f"Error reading file {file_name}: {e}")
        return f"Error reading file {file_name}: {e}"

async def clean_html(page: Page) -> str:
    """
    清理 HTML 中的脚本和样式标签
    """

    # 用evaluate获取html
    html = ""
    useful_elements = []
    # 需要去重
    elements = page.query_selector_all(":not(:has(*))")
    for element in elements:
        outer_html = await element.evaluate("(el) => el.outerHTML")
        # outer_html = str(element)
        # print(outer_html)
        if outer_html not in useful_elements:
            useful_elements.append(outer_html)
            # print(outer_html)

    for outer_html in useful_elements:
        html += outer_html + "\n\n"

    with open("crewler1.html", "w", encoding="utf-8") as f:
        f.write(html)

    soup = BeautifulSoup(html, "html.parser")
    soup.prettify()

    # 移除 HTML 注释
    for comment in soup.find_all(string=lambda text: isinstance(text, Comment)):
        comment.extract()

    # 需要保留 script 标签，因为可能包含有效信息
    for script in soup(["style", "meta", "link", "path", "noscript"]):
        script.decompose()

    interactive_elements = [
        "a",  # 链接
        "select",  # 下拉框
        "img",  # 图片
        "[href]:not(a)",  # 具有链接的元素
    ]

    for element in soup.find_all(True):
        if element.name not in interactive_elements: # type:ignore
            element.unwrap() # type:ignore

    # 去除所有空行和重复行
    html = str(soup)
    # 清理超过指定长度的行
    html = clean_overlength_lines(html)
    html = re.sub(r"\n\s*\n", "\n", html)
    lines = html.split("\n")
    seen = set()
    unique_lines = []
    for line in lines:
        line = line.strip()
        if line and line not in seen:
            seen.add(line)
            unique_lines.append(line)
    html = "\n".join(unique_lines)

    with open("crewler.html", "w", encoding="utf-8") as f:
        f.write(html)

    return html


def clean_overlength_lines(content: str) -> str:
    """
    清理超过指定长度的行
    """
    lines = content.split("\n")
    for line in lines:
        # 如果行长度超过指定长度，则删除该行
        if len(line) > LINE_LIMIT:
            # 删除该行
            lines.remove(line)
    return "\n".join(lines)


def get_last_n_chars(text: str, n: int = 10000) -> str:
    """
    从字符串末尾获取指定数量的字符
    Args:
        text: 输入字符串
        n: 需要获取的字符数量，默认为1000
    Returns:
        str: 从末尾获取的n个字符，如果发生截断会包含长度信息
    """
    if not text:
        return ""
    if len(text) > n:
        return f"<!-- Original length: {len(text)}, truncated to: {n} -->\n{text[-n:]}"
    return text


def get_head_and_tail_n_chars(text: str, n: int = 10000) -> str:
    """
    获取文本的头部和尾部指定数量的字符
    """
    if not text:
        return ""
    if len(text) > n:
        half_n = n // 2
        return f"<!-- Original length: {len(text)} bytes, truncated to: head {half_n} bytes and tail {half_n} bytes -->\n{text[:half_n]}\n...\n{text[-half_n:]}"
    return text


def file_to_text_with_line_number(file_path: str) -> str:
    """
    将文件转换为带有行号的文本
    """
    text = ""
    with open(file_path, "r", encoding="utf-8") as f:
        lines = f.readlines()
        text += "".join([f"{i+1}: {line}" for i, line in enumerate(lines)])
    return text


def clean_tqdm_output(output: str) -> str:
    """Clean tqdm progress bars from output."""
    if not output:
        return ""

    # Remove tqdm progress bars (they typically contain \r and %)
    lines = output.split("\n")
    cleaned_lines = []

    for line in lines:
        # Skip tqdm progress bars
        if "\r" in line and ("%" in line or "it/s" in line):
            continue
        # Skip lines that are just progress indicators
        if line.strip().endswith("%") or line.strip().endswith("it/s"):
            continue
        # Skip lines that are just progress bars
        if re.match(r"^\s*[\d.]+\%|\d+/\d+|\d+it/s", line.strip()):
            continue
        cleaned_lines.append(line)

    return "\n".join(cleaned_lines)


def get_file_segment(
    dir: str, file_name: str, line_numbers: List[Dict[str, int]]
) -> str:
    """
    获取文件的指定行段
    """
    file_path = os.path.join(dir, file_name)
    content = ""
    # 每一行前面加上行号
    for line_number in line_numbers:
        content += f"<!-- Line {line_number['start']} to {line_number['end']} -->\n"
        with open(file_path, "r", encoding="utf-8") as f:
            lines = f.readlines()[line_number["start"] - 1 : line_number["end"]]
            start_line_number = line_number["start"]
            end_line_number = line_number["end"]
            content += "".join(
                [f"{i+start_line_number}: {line}" for i, line in enumerate(lines)]
            )
    return content

async def wait_for_confirmation_input(caller) -> str:
    """
    等待用户确认(通过输入)
    返回值（str）:
        "true": 确认
        "reExecute": 重新执行
        "index:数字": 返回到某次执行
        "message:消息内容": 穿插用户消息
    """
    print("\n请选择操作:")
    print("1. 确认继续 (输入 y)")
    print("2. 重新执行 (输入 r)")
    print("3. 返回到指定执行点 (输入 i)")
    print("4. 插入消息 (输入 m)")

    index = int(caller["index"])

    while True:
        choice = input("\n请输入您的选择: ").strip().lower()
        
        if choice == 'y':
            return "true"
        elif choice == 'r':
            backtrack_task_execute_history(caller=caller, index=index)
            return "reExecute"
        elif choice == 'i':
            while True:
                index = input("请输入要返回的执行点编号: ").strip()
                if index.isdigit():
                    backtrack_task_execute_history(caller, int(index))
                    return f"index"
                print("请输入有效的数字!")
        elif choice == 'm':
            message = input("请输入要插入的消息: ").strip()
            write_user_message(caller, message)
            if message:
                return f"message"
            print("消息不能为空!")
        else:
            print("无效的选择,请重新输入!")