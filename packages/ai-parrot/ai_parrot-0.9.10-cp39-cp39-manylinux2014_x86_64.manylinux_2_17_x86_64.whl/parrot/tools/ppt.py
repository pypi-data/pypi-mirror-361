import re
import uuid
from pathlib import Path
from typing import Dict, List, Optional, Type, Any
from datetime import datetime
import asyncio
import traceback
from langchain.tools import BaseTool
from pydantic import BaseModel, Field, ConfigDict
from jinja2 import Environment, FileSystemLoader
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import markdown
from bs4 import BeautifulSoup, NavigableString
from navconfig import BASE_DIR
from .abstract import BaseAbstractTool


class PowerPointInput(BaseModel):
    """
    Input schema for the PowerPointGeneratorTool. Users can supply:
    • text (required): the markdown text to convert to PowerPoint slides
    • output_filename: (Optional) custom filename (including extension) for the generated presentation
    • template_name: (Optional) name of the HTML template (e.g. 'presentation.html') to render
    • template_vars: (Optional) dict of variables to pass into the template
    • output_dir: (Optional) directory where the PowerPoint file will be saved
    • pptx_template: (Optional) path to a PowerPoint template file to use as base
    • slide_layout: (Optional) default slide layout index (default: 1)
    • title_styles: (Optional) dict of styles to apply to slide titles
    • content_styles: (Optional) dict of styles to apply to slide content
    """
    model_config = ConfigDict(extra='forbid')

    text: str = Field(..., description="The markdown text to convert to PowerPoint slides")
    output_filename: Optional[str] = Field(
        None,
        description="(Optional) Custom filename (including extension) for the generated presentation"
    )
    template_name: Optional[str] = Field(
        None,
        description="Name of the HTML template (e.g. 'presentation.html') to render"
    )
    template_vars: Optional[Dict[str, Any]] = Field(
        None,
        description="Dict of variables to pass into the template"
    )
    output_dir: Optional[str] = Field(
        None,
        description="Directory where the PowerPoint file will be saved"
    )
    pptx_template: Optional[str] = Field(
        None,
        description="Path to PowerPoint template file to use as base"
    )
    slide_layout: Optional[int] = Field(
        1,
        description="Default slide layout index (0=Title Slide, 1=Title and Content, etc.)"
    )
    title_styles: Optional[Dict[str, Any]] = Field(
        None,
        description="Dict of styles to apply to slide titles (font, size, color, etc.)"
    )
    content_styles: Optional[Dict[str, Any]] = Field(
        None,
        description="Dict of styles to apply to slide content"
    )


class PowerPointGeneratorTool(BaseAbstractTool):
    """PowerPoint Presentation Generator Tool.

    * How Slide Splitting Works:

    # Main Title          → Slide 1 (Title slide)
    Content here...

    ## Section 1          → Slide 2 (Title + Content)
    Paragraphs and lists...

    ### Subsection       → Slide 3 (Title + Content)
    More content...

    ## Section 2          → Slide 4 (Title + Content)
    Tables and text...
    """
    name: str = "generate_powerpoint_presentation"
    description: str = (
        "Create PowerPoint presentations from markdown text, "
        "splitting content into slides based on headings."
    )
    env: Optional[Environment] = None
    templates_dir: Optional[Path] = None

    args_schema: Type[BaseModel] = PowerPointInput

    def __init__(
        self,
        *args,
        templates_dir: Optional[Path] = None,
        **kwargs
    ):
        """Initialize the PowerPoint generator tool."""
        super().__init__(*args, **kwargs)
        # Initialize Jinja2 environment for HTML templates
        if templates_dir:
            self.templates_dir = templates_dir
            self.env = Environment(
                loader=FileSystemLoader(str(templates_dir)),
                autoescape=True
            )

    def _default_output_dir(self) -> Path:
        """Get default output directory for PDF files."""
        return self.static_dir.joinpath('documents', 'presentations')

    def _generate_payload(
        self,
        **kwargs
    ) -> PowerPointInput:
        try:
            # Validate input using Pydantic
            return PowerPointInput(**kwargs)
        except Exception as e:
            raise ValueError(f"Invalid input for PowerPoint generation: {e}")

    async def _generate_content(self, input_data: PowerPointInput) -> Dict[str, Any]:
        """Generate PowerPoint presentation from markdown text."""
        try:
            # Process the text through Jinja2 template if provided
            processed_text = self._render_template(input_data)

            # Preprocess markdown
            processed_text = self._preprocess_markdown(processed_text)

            # Convert markdown to HTML
            html_content = self._markdown_to_html(processed_text)

            # Parse HTML and extract slides
            slides_data = self._extract_slides_from_html(html_content)

            # Create or load PowerPoint presentation
            prs = self._create_presentation(input_data.pptx_template)

            # Create slides from extracted data
            self._create_slides(prs, slides_data, input_data)

            # Save presentation
            output_path = self._save_presentation(prs, input_data)
            url = self.to_static_url(output_path)
            return {
                "status": "success",
                "file_path": output_path,
                "filename": output_path,
                "slides_created": len(slides_data),
                "url": url,
                "static_url": self.relative_url(url),
                "output_filename": input_data.output_filename or output_path.name,
                "file_path": self.output_dir,
                "message": f"PowerPoint presentation successfully created at {output_path}"
            }

        except Exception as e:
            return {
                "status": "error",
                "error": str(e),
                "message": f"Failed to create PowerPoint presentation: {str(e)}"
            }

    def _render_template(self, input_data: PowerPointInput) -> str:
        """Render text through Jinja2 template if provided."""
        if not input_data.template_name or not self.env:
            return input_data.text

        try:
            template = self.env.get_template(input_data.template_name)
            template_vars = input_data.template_vars or {}

            # Add some default variables
            template_vars.setdefault('content', input_data.text)
            template_vars.setdefault('date', datetime.now().strftime('%Y-%m-%d'))
            template_vars.setdefault('timestamp', datetime.now().isoformat())

            return template.render(**template_vars)
        except Exception as e:
            print(f"Template rendering failed: {e}")
            return input_data.text

    def _preprocess_markdown(self, text: str) -> str:
        """Preprocess markdown to handle common issues."""
        # Replace placeholder variables with empty strings
        text = re.sub(r'\{[a-zA-Z0-9_]+\}', '', text)

        # Handle f-strings that weren't evaluated
        text = re.sub(r'f"""(.*?)"""', r'\1', text, flags=re.DOTALL)
        text = re.sub(r"f'''(.*?)'''", r'\1', text, flags=re.DOTALL)

        # Remove triple backticks and language indicators
        text = re.sub(r'```[a-zA-Z]*\n', '', text)
        text = re.sub(r'```', '', text)

        # Fix heading issues (ensure space after #)
        text = re.sub(r'(#+)([^ \n])', r'\1 \2', text)

        return text

    def _markdown_to_html(self, markdown_text: str) -> str:
        """Convert markdown to HTML."""
        try:
            html = markdown.markdown(
                markdown_text,
                extensions=['extra', 'codehilite', 'toc', 'tables']
            )
            return html
        except Exception as e:
            print(f"Markdown conversion failed: {e}")
            # Fallback: wrap in paragraphs
            paragraphs = markdown_text.split('\n\n')
            html_paragraphs = [f'<p>{p.replace(chr(10), "<br>")}</p>' for p in paragraphs if p.strip()]
            return '\n'.join(html_paragraphs)

    def _extract_slides_from_html(self, html_content: str) -> List[Dict[str, Any]]:
        """Extract slides from HTML content based on headings."""
        soup = BeautifulSoup(html_content, 'html.parser')
        slides = []

        # Find all heading elements
        headings = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])

        if not headings:
            # If no headings, create a single slide with all content
            slides.append({
                'title': 'Slide',
                'content': self._extract_content_elements(soup),
                'level': 1
            })
            return slides

        for i, heading in enumerate(headings):
            slide_data = {
                'title': self._get_text_content(heading),
                'level': int(heading.name[1]),  # h1 -> 1, h2 -> 2, etc.
                'content': []
            }

            # Find content between this heading and the next one
            current_element = heading.next_sibling
            next_heading = headings[i + 1] if i + 1 < len(headings) else None

            while current_element and current_element != next_heading:
                if hasattr(current_element, 'name') and current_element.name in ['p', 'ul', 'ol', 'table', 'blockquote', 'div']:
                    slide_data['content'].append(current_element)
                current_element = current_element.next_sibling

            slides.append(slide_data)

        return slides

    def _extract_content_elements(self, soup) -> List:
        """Extract content elements from soup."""
        content_elements = []
        for element in soup.find_all(['p', 'ul', 'ol', 'table', 'blockquote', 'div']):
            content_elements.append(element)
        return content_elements

    def _create_presentation(self, template_path: Optional[str] = None) -> Presentation:
        """Create or load PowerPoint presentation."""
        if template_path and Path(template_path).exists():
            return Presentation(template_path)
        else:
            return Presentation()

    def _create_slides(self, prs: Presentation, slides_data: List[Dict[str, Any]], input_data: PowerPointInput):
        """Create slides from extracted data."""
        for i, slide_data in enumerate(slides_data):
            # Determine slide layout based on heading level or use default
            if i == 0 and slide_data['level'] == 1:
                # First slide with h1 might be title slide
                layout_idx = 0  # Title slide
            else:
                layout_idx = input_data.slide_layout or 1  # Title and content

            # Add slide
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            # Add title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
                if input_data.title_styles:
                    self._apply_text_styles(slide.shapes.title, input_data.title_styles)

            # Add content
            if slide_data['content'] and len(slide.shapes.placeholders) > 1:
                content_placeholder = slide.shapes.placeholders[1]

                # Clear existing content
                content_placeholder.text = ""

                # Add content based on type
                self._add_slide_content(content_placeholder, slide_data['content'], input_data)

    def _add_slide_content(self, placeholder, content_elements: List, input_data: PowerPointInput):
        """Add content to a slide placeholder."""
        text_frame = placeholder.text_frame
        text_frame.clear()

        for element in content_elements:
            if element.name == 'p':
                # Add paragraph
                p = text_frame.paragraphs[0] if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text else text_frame.add_paragraph()
                p.text = self._get_text_content(element)
                if input_data.content_styles:
                    self._apply_paragraph_styles(p, input_data.content_styles)

            elif element.name in ['ul', 'ol']:
                # Add list items
                for li in element.find_all('li', recursive=False):
                    p = text_frame.add_paragraph()
                    p.text = self._get_text_content(li)
                    p.level = 1  # Bullet point level
                    if input_data.content_styles:
                        self._apply_paragraph_styles(p, input_data.content_styles)

            elif element.name == 'table':
                # Add table (simplified - would need more work for complex tables)
                table_text = self._extract_table_text(element)
                p = text_frame.add_paragraph()
                p.text = table_text
                if input_data.content_styles:
                    self._apply_paragraph_styles(p, input_data.content_styles)

            elif element.name == 'blockquote':
                # Add blockquote
                p = text_frame.add_paragraph()
                p.text = f'"{self._get_text_content(element)}"'
                if input_data.content_styles:
                    self._apply_paragraph_styles(p, input_data.content_styles)

    def _extract_table_text(self, table_element) -> str:
        """Extract text from table element."""
        rows = table_element.find_all('tr')
        table_lines = []

        for row in rows:
            cells = row.find_all(['td', 'th'])
            row_text = ' | '.join([self._get_text_content(cell) for cell in cells])
            table_lines.append(row_text)

        return '\n'.join(table_lines)

    def _get_text_content(self, element) -> str:
        """Extract text content from HTML element."""
        if isinstance(element, NavigableString):
            return str(element)

        text_parts = []
        for content in element.contents:
            if isinstance(content, NavigableString):
                text_parts.append(str(content))
            else:
                text_parts.append(self._get_text_content(content))

        return ''.join(text_parts).strip()

    def _apply_text_styles(self, shape, styles: Dict[str, Any]):
        """Apply styles to a text shape."""
        if not shape.has_text_frame:
            return

        text_frame = shape.text_frame

        for paragraph in text_frame.paragraphs:
            self._apply_paragraph_styles(paragraph, styles)

    def _apply_paragraph_styles(self, paragraph, styles: Dict[str, Any]):
        """Apply styles to a paragraph."""
        # Font styling
        if 'font_name' in styles:
            paragraph.font.name = styles['font_name']
        if 'font_size' in styles:
            paragraph.font.size = Pt(styles['font_size'])
        if 'bold' in styles:
            paragraph.font.bold = styles['bold']
        if 'italic' in styles:
            paragraph.font.italic = styles['italic']
        if 'font_color' in styles:
            # Convert hex color to RGB
            color_hex = styles['font_color'].lstrip('#')
            r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
            paragraph.font.color.rgb = RGBColor(r, g, b)

        # Alignment
        if 'alignment' in styles:
            alignment_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT,
                'justify': PP_ALIGN.JUSTIFY
            }
            paragraph.alignment = alignment_map.get(styles['alignment'], PP_ALIGN.LEFT)

    def _save_presentation(self, prs: Presentation, input_data: PowerPointInput) -> Path:
        """Save the presentation and return the file path."""
        # Determine output directory
        output_dir = Path(input_data.output_dir) if input_data.output_dir else self.output_dir
        output_dir.mkdir(parents=True, exist_ok=True)

        if input_data.output_filename:
            output_filename = input_data.output_filename
            if not output_filename.endswith('.pptx'):
                output_filename += '.pptx'
        else:
            output_filename = self.generate_filename(
                prefix=input_data.file_prefix or "presentation",
                extension="pptx"
            )
        output_path = output_dir / output_filename

        # Save the presentation
        prs.save(str(output_path))

        return output_path
